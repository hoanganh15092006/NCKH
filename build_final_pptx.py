import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

def clear_slides(prs):
    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    for slide in slides:
        xml_slides.remove(slide)

def add_slide_with_image(prs, layout_idx, title, bullets, img_path=None):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    
    if slide.shapes.title:
        slide.shapes.title.text = title
        
    # Attempt to find body placeholder
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
            
    if body and bullets:
        tf = body.text_frame
        tf.clear()
        for i, text in enumerate(bullets):
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(24)
            
    if img_path:
        # Add picture on the right side
        left = Inches(5.0)
        top = Inches(1.5)
        height = Inches(5)
        try:
            slide.shapes.add_picture(img_path, left, top, height=height)
        except Exception as e:
            print(f"Lỗi chèn ảnh: {e}")

prs = Presentation(r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\slide NCKH.pptx")
clear_slides(prs)

# Add Slide 1: Title
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
if title_slide.shapes.title:
    title_slide.shapes.title.text = "HỆ THỐNG QUẢN LÝ BÃI ĐỖ XE THÔNG MINH\nSMART PARKING"
try:
    title_slide.placeholders[1].text = "Sinh viên thực hiện: Hoàng Anh\nĐại học Giao thông Vận tải"
except:
    pass

# Hình ảnh đã tạo bằng AI
img_gate = r"C:\Users\Admin\.gemini\antigravity\brain\51bc3eeb-dbfc-4d66-acdd-8cc35691ea1e\smart_parking_gate_1777994554269.png"
img_app = r"C:\Users\Admin\.gemini\antigravity\brain\51bc3eeb-dbfc-4d66-acdd-8cc35691ea1e\parking_mobile_app_1777994670755.png"
img_ai = r"C:\Users\Admin\.gemini\antigravity\brain\51bc3eeb-dbfc-4d66-acdd-8cc35691ea1e\ai_neural_network_1777994932934.png"

# Các slide nội dung
add_slide_with_image(prs, 1, "1. Thực trạng các bãi đỗ xe hiện nay", ["- Ùn tắc cục bộ vào giờ cao điểm", "- Lấy thẻ, trả tiền mặt kéo dài thời gian", "- Khói bụi và tiếng ồn tại khu vực soát vé"], img_gate)
add_slide_with_image(prs, 1, "2. Rủi ro của phương pháp truyền thống", ["- Thẻ từ vật lý dễ bị mất mát, đánh tráo", "- Vé giấy dễ làm giả hoặc rách nát do mưa", "- Chi phí nhân sự túc trực 24/7 cực kỳ lớn"])
add_slide_with_image(prs, 1, "3. Mục tiêu của Dự án NCKH", ["- Hiện đại hóa hoàn toàn quy trình gửi xe.", "- Tự động hóa thanh toán, tiến tới xã hội không tiền mặt.", "- Ứng dụng công nghệ AI để tối ưu trải nghiệm người dùng."])
add_slide_with_image(prs, 1, "4. Giải pháp: Triết lý 3 KHÔNG", ["- KHÔNG CHỜ ĐỢI: Camera nhận diện AI siêu tốc.", "- KHÔNG CHẠM: Barie điều khiển mở tự động.", "- KHÔNG TIỀN MẶT: Thanh toán tự động qua App."])
add_slide_with_image(prs, 1, "5. Kiến trúc Hệ thống Tổng thể", ["Hệ thống phân tán với 3 cụm liên kết chặt chẽ:", "1. Cụm AI Local: Mắt thần đọc biển số xe.", "2. Backend Server: Bộ não Python Flask/MySQL.", "3. Mobile App: Android Retrofit làm điểm chạm."])
add_slide_with_image(prs, 1, "6. Lõi công nghệ 1: Camera AI Nhận diện", ["- Hoạt động liên tục tại cửa vào/ra.", "- Sử dụng Computer Vision và Deep Learning.", "- Bóc tách ký tự và gửi lệnh trong vài mili-giây."], img_gate)
add_slide_with_image(prs, 1, "7. Lõi công nghệ 2: Backend Server", ["- Bộ trung tâm xử lý dữ liệu và điều phối.", "- Cung cấp RESTful API tốc độ cực cao.", "- Giao tiếp thời gian thực với phần cứng Barie.", "- Lưu trữ dữ liệu an toàn trên MySQL."])
add_slide_with_image(prs, 1, "8. Trải nghiệm: Mobile App", ["- Giao diện tối giản, hiện đại (Tông xanh Navy).", "- Quản lý phương tiện và theo dõi trạng thái Real-time.", "- Quét QR dự phòng khi Camera bị chói lóa/lỗi.", "- Quản lý số dư ví điện tử cá nhân."], img_app)
add_slide_with_image(prs, 1, "9. Thách thức lớn nhất: Thanh toán", ["- Bài toán: Làm sao để nạp tiền tự động 24/7?", "- Các hệ thống cũ yêu cầu Admin duyệt tay (chậm chạp).", "- Hoặc trả phí đắt đỏ cho VNPAY, MoMo.", "=> Cần giải pháp đối soát ngân hàng Miễn phí & Tự động."])
add_slide_with_image(prs, 1, "10. Giải pháp: Robot tự động RPA", ["- Phát triển công cụ RPA (Robotic Process Automation).", "- Sử dụng NodeJS/Puppeteer chạy nền trên máy chủ.", "- Tự động đăng nhập Internet Banking của MB Bank.", "- Quét lịch sử giao dịch liên tục mỗi 5 giây."])
add_slide_with_image(prs, 1, "11. Rào cản kỹ thuật: CAPTCHA Ngân hàng", ["- Ngân hàng MB Bank dùng ảnh CAPTCHA để chống Robot.", "- Các tool RPA bình thường sẽ bị chặn 100%.", "- Giải pháp: Bắt buộc phải có trí tuệ nhân tạo can thiệp."])
add_slide_with_image(prs, 1, "12. Đột phá: AI Giải mã CAPTCHA", ["- Tự huấn luyện mô hình Mạng Nơ-ron Tích chập (CNN).", "- Sử dụng nền tảng TensorFlow mạnh mẽ.", "- AI tự động khử nhiễu, tách nền và đọc chữ số.", "- Hiệu suất giải mã chính xác cao, thời gian tính bằng ms."], img_ai)
add_slide_with_image(prs, 1, "13. Quy trình Thanh toán Khép kín", ["- B1: App sinh mã VietQR chuyên biệt.", "- B2: Khách chuyển khoản -> MB Bank nhận tiền.", "- B3: AI vượt rào, Robot báo về Server.", "- B4: Server cộng tiền vào ví trong 3 giây. KHÔNG CẦN CON NGƯỜI."])
add_slide_with_image(prs, 1, "14. Đánh giá Kết quả Đạt được", ["- Xây dựng thành công toàn bộ hệ sinh thái phần mềm.", "- Thời gian lưu thông qua cổng cực ngắn (< 3 giây).", "- Ứng dụng thành công AI vào thực tiễn phức tạp (ngân hàng)."])
add_slide_with_image(prs, 1, "15. Định hướng Mở rộng trong tương lai", ["- Cải tiến thuật toán chỉ đường (Indoor Navigation) cho bãi đỗ.", "- Mở rộng App trên nền tảng iOS (Apple).", "- Hỗ trợ đa dạng ngân hàng hơn ngoài MB Bank."])
add_slide_with_image(prs, 0, "XIN CHÂN THÀNH CẢM ƠN HỘI ĐỒNG!", ["Em rất mong nhận được những góp ý, phản biện từ quý Thầy Cô", "để đề tài được hoàn thiện hơn."])

prs.save(r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\Slide_NCKH_HoanThien.pptx")
print("Hoàn thành!")
