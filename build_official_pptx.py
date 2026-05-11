import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def clear_slides(prs):
    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    for slide in slides:
        xml_slides.remove(slide)

def add_slide_with_image(prs, layout_idx, title, bullets, notes="", img_path=None):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    
    if slide.shapes.title:
        slide.shapes.title.text = title
        try:
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)
        except: pass
        
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
            
    if body and bullets:
        tf = body.text_frame
        tf.clear()
        for i, text in enumerate(bullets):
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(64, 64, 64)
            
    if img_path and os.path.exists(img_path):
        left = Inches(5.2)
        top = Inches(2.0)
        height = Inches(4.0)
        try:
            slide.shapes.add_picture(img_path, left, top, height=height)
        except Exception as e:
            print(f"Lỗi chèn ảnh: {e}")

    if notes and slide.has_notes_slide:
        slide.notes_slide.notes_text_frame.text = notes

prs = Presentation(r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\slide NCKH.pptx")
clear_slides(prs)

# SLIDE 1
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
if title_slide.shapes.title:
    title_slide.shapes.title.text = "Ứng dụng trí tuệ nhân tạo trong nhận diện biển số xe\nvà xây dựng hệ thống quản lý bãi đỗ xe thông minh"
    try: title_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)
    except: pass
try:
    title_slide.placeholders[1].text = "Trường Đại học Giao thông Vận tải\nĐề tài Nghiên cứu khoa học sinh viên năm 2026\nSinh viên thực hiện: ...\nGiảng viên hướng dẫn: TS. Phạm Đình Phong"
except: pass
if title_slide.has_notes_slide:
    title_slide.notes_slide.notes_text_frame.text = "Kính thưa quý thầy cô, em xin đại diện nhóm trình bày đề tài nghiên cứu khoa học sinh viên: Ứng dụng trí tuệ nhân tạo trong nhận diện biển số xe và xây dựng hệ thống quản lý bãi đỗ xe thông minh."

# Image mappings (fallback to web_*.jpg if the chromium ones fail)
img_traffic = r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\traffic.jpg"
if not os.path.exists(img_traffic): img_traffic = r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\web_parking.jpg"

img_yolo = r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\yolo.jpg"
if not os.path.exists(img_yolo): img_yolo = r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\web_ai.jpg"

img_app = r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\web_app.jpg"


slides_data = [
    {
        "title": "NỘI DUNG TRÌNH BÀY",
        "bullets": ["- Bài toán thực tế và mục tiêu đề tài", "- Kiến trúc hệ thống và công nghệ sử dụng", "- Module AI nhận diện biển số", "- Ứng dụng hệ thống và demo sản phẩm", "- Kết luận và hướng phát triển"],
        "notes": "Bài trình bày của nhóm em gồm 5 phần chính: đầu tiên là bài toán thực tế và mục tiêu, sau đó là kiến trúc hệ thống, phần AI nhận diện biển số, phần ứng dụng và demo, cuối cùng là kết luận và hướng phát triển.",
        "img": None
    },
    {
        "title": "BÀI TOÁN THỰC TẾ",
        "bullets": ["- Số lượng phương tiện cá nhân ngày càng tăng", "- Nhu cầu quản lý bãi đỗ xe ngày càng lớn", "- Nhiều bãi xe vẫn quản lý thủ công hoặc bán tự động", "- Nhân viên phải kiểm tra thẻ, ghi biển số hoặc nhập dữ liệu bằng tay"],
        "notes": "Trong thực tế, số lượng phương tiện cá nhân ngày càng tăng, đặc biệt ở các khu vực đô thị, trường học, chung cư và trung tâm thương mại. Điều này khiến nhu cầu quản lý bãi đỗ xe trở nên cấp thiết hơn. Tuy nhiên, nhiều bãi xe hiện nay vẫn còn phụ thuộc vào thao tác thủ công như ghi biển số, kiểm tra thẻ hoặc nhập dữ liệu bằng tay.",
        "img": img_traffic
    },
    {
        "title": "HẠN CHẾ CỦA CÁCH QUẢN LÝ TRUYỀN THỐNG",
        "bullets": ["- Xử lý xe vào/ra chậm", "- Dễ sai sót khi nhập biển số", "- Khó tra cứu lịch sử giao dịch", "- Tốn nhân lực vận hành", "- Người dùng khó tự kiểm tra thông tin"],
        "notes": "Cách quản lý truyền thống có một số hạn chế. Thứ nhất là tốc độ xử lý xe vào ra chậm, đặc biệt khi lượng xe đông. Thứ hai là dễ sai sót khi nhân viên nhập biển số hoặc thời gian thủ công. Thứ ba là việc tra cứu lịch sử không thuận tiện. Ngoài ra, người dùng cũng khó tự kiểm tra các thông tin như số dư hoặc lịch sử gửi xe.",
        "img": img_traffic
    },
    {
        "title": "MỤC TIÊU ĐỀ TÀI",
        "bullets": ["- Xây dựng hệ thống quản lý bãi đỗ xe thông minh", "- Nhận diện biển số xe tự động bằng AI", "- Xây dựng phần mềm quản lý cho nhân viên", "- Xây dựng ứng dụng di động cho người dùng", "- Kết nối hệ thống qua API và cơ sở dữ liệu"],
        "notes": "Từ những vấn đề trên, nhóm em đặt ra mục tiêu xây dựng một hệ thống quản lý bãi đỗ xe thông minh. Hệ thống tích hợp công nghệ nhận diện biển số tự động bằng AI, đồng thời có phần mềm quản lý cho nhân viên và ứng dụng di động cho người dùng. Các thành phần được kết nối với nhau thông qua API và cơ sở dữ liệu.",
        "img": None
    },
    {
        "title": "PHẠM VI NGHIÊN CỨU",
        "bullets": ["- Đối tượng: biển số xe Việt Nam", "- Công nghệ nhận diện: YOLOv8 + EasyOCR", "- Phần mềm quản lý: Desktop App", "- Ứng dụng người dùng: Android App", "- Cơ sở dữ liệu và API kết nối hệ thống"],
        "notes": "Trong phạm vi nghiên cứu, nhóm em tập trung vào nhận diện biển số xe Việt Nam tự động. Hệ thống sử dụng YOLOv8 để phát hiện vùng biển số và EasyOCR để đọc ký tự. Bên cạnh đó, nhóm xây dựng phần mềm quản lý trên máy tính, ứng dụng Android cho người dùng, cùng cơ sở dữ liệu và API để kết nối các thành phần.",
        "img": img_yolo
    },
    {
        "title": "KIẾN TRÚC TỔNG THỂ HỆ THỐNG",
        "bullets": ["1. Camera / Video", "      ↓", "2. YOLOv8 phát hiện biển số", "      ↓", "3. EasyOCR nhận dạng ký tự", "      ↓", "4. Backend / API + Database", "      ↓", "5. Desktop App / Android App"],
        "notes": "Về kiến trúc tổng thể, hệ thống bắt đầu từ camera hoặc video đầu vào. Hình ảnh được đưa vào YOLOv8 để phát hiện vùng biển số, sau đó vùng biển số được cắt ra và đưa vào EasyOCR để nhận dạng ký tự. Kết quả nhận diện được lưu vào cơ sở dữ liệu thông qua backend/API, sau đó hiển thị trên phần mềm quản lý và ứng dụng người dùng.",
        "img": None
    },
    {
        "title": "LUỒNG XE VÀO / XE RA",
        "bullets": ["* Xe vào:", "- Camera → nhận diện biển số → tạo phiên gửi xe → lưu thời gian vào", "* Xe ra:", "- Camera → nhận diện lại biển số → đối chiếu phiên gửi xe → cập nhật lịch sử", "=> Biển số là thông tin định danh chính"],
        "notes": "Đối với xe vào, camera ghi lại hình ảnh xe, hệ thống nhận diện biển số và tạo một phiên gửi xe mới. Đối với xe ra, hệ thống nhận diện lại biển số, đối chiếu với phiên gửi xe đang hoạt động, sau đó cập nhật lịch sử giao dịch. Trong hệ thống này, biển số đóng vai trò là thông tin định danh chính của phương tiện.",
        "img": None
    },
    {
        "title": "CÔNG NGHỆ SỬ DỤNG",
        "bullets": ["- Python: xử lý ảnh, AI, phần mềm quản lý", "- OpenCV: đọc camera/video và xử lý ảnh", "- YOLOv8: phát hiện vùng biển số", "- EasyOCR: nhận dạng ký tự", "- Flask REST API: kết nối app và server", "- Database: lưu user, biển số, phiên gửi xe, lịch sử", "- Android: ứng dụng người dùng"],
        "notes": "Về công nghệ, nhóm em sử dụng Python cho phần xử lý ảnh, AI và phần mềm quản lý. OpenCV dùng để đọc camera hoặc video. YOLOv8 dùng để phát hiện vùng biển số, còn EasyOCR dùng để đọc ký tự. Flask REST API đóng vai trò server trung gian, giúp ứng dụng Android giao tiếp với cơ sở dữ liệu.",
        "img": img_app
    },
    {
        "title": "CƠ SỞ DỮ LIỆU VÀ API",
        "bullets": ["* Cơ sở dữ liệu lưu:", "- USERS: tài khoản người dùng", "- OWNED_PLATES: biển số của người dùng", "- ACTIVE_SESSIONS: phiên gửi xe đang hoạt động", "- HISTORY_RECORDS: lịch sử giao dịch", "- SETTINGS: cài đặt hệ thống", "- REMOTE_COMMANDS: lệnh điều khiển từ xa", "=> API trao đổi dữ liệu giữa Android App và hệ thống"],
        "notes": "Cơ sở dữ liệu của hệ thống lưu các nhóm thông tin chính như tài khoản người dùng, biển số đã đăng ký, phiên gửi xe đang hoạt động và lịch sử giao dịch. Backend API giúp ứng dụng Android gửi yêu cầu như đăng nhập, xem số dư hoặc xem lịch sử, sau đó server truy vấn dữ liệu và trả kết quả về cho app.",
        "img": None
    },
    {
        "title": "PIPELINE AI NHẬN DIỆN BIỂN SỐ",
        "bullets": ["1. Ảnh / Video đầu vào", "2. → Phát hiện vùng biển số", "3. → Cắt ảnh biển số", "4. → Tiền xử lý ảnh", "5. → OCR đọc ký tự", "6. → Hậu xử lý", "7. → Kết quả biển số"],
        "notes": "Pipeline nhận diện biển số của nhóm gồm nhiều bước. Đầu tiên là lấy ảnh hoặc video đầu vào. Sau đó YOLOv8 phát hiện vùng biển số. Vùng này được cắt ra và tiền xử lý để tăng chất lượng ảnh. Tiếp theo, EasyOCR đọc ký tự và kết quả được hậu xử lý để đưa ra biển số cuối cùng.",
        "img": img_yolo
    },
    {
        "title": "YOLOv8 PHÁT HIỆN VÙNG BIỂN SỐ",
        "bullets": ["- Sử dụng mô hình YOLOv8s", "- Đầu vào: frame ảnh/video", "- Đầu ra: bounding box vùng biển số", "* Ưu điểm:", "- Tốc độ nhanh", "- Phù hợp xử lý gần thời gian thực", "- Có thể fine-tune cho bài toán biển số xe"],
        "notes": "Trong hệ thống, YOLOv8s được dùng để phát hiện vị trí biển số trong khung hình. Mô hình không đọc nội dung biển số, mà chỉ xác định vùng biển số nằm ở đâu. Kết quả đầu ra là bounding box bao quanh biển số, sau đó vùng ảnh này sẽ được đưa sang bước OCR.",
        "img": img_yolo
    },
    {
        "title": "DỮ LIỆU VÀ HUẤN LUYỆN MÔ HÌNH",
        "bullets": ["- Dataset: Roboflow Universe (Định dạng YOLOv8)", "- Tập train: 6906 ảnh", "- Tập validation: 1973 ảnh", "- Model: YOLOv8s pretrained", "- Huấn luyện: 50 epoch", "- Môi trường: Google Colab, GPU Tesla T4", "- Trọng số sử dụng: best.pt"],
        "notes": "Về dữ liệu huấn luyện, nhóm em sử dụng tập dữ liệu biển số từ Roboflow Universe và tải về theo định dạng YOLOv8. Tập dữ liệu gồm khoảng 6906 ảnh train và 1973 ảnh validation. Nhóm không huấn luyện mô hình từ đầu, mà sử dụng YOLOv8s pretrained làm trọng số khởi tạo, sau đó fine-tune trong 50 epoch trên Google Colab có GPU Tesla T4. Sau huấn luyện, file best.pt được dùng trong hệ thống nhận diện.",
        "img": img_yolo
    },
    {
        "title": "EASYOCR, TIỀN XỬ LÝ VÀ HẬU XỬ LÝ",
        "bullets": ["* Tiền xử lý ảnh:", "- Cắt vùng biển số, Chuyển ảnh xám, Tăng tương phản, Khử nhiễu, Làm rõ ký tự", "* EasyOCR đọc ký tự trên biển số", "* Hậu xử lý lỗi OCR:", "- O ↔ 0", "- I ↔ 1", "- S ↔ 5", "- B ↔ 8"],
        "notes": "Sau khi YOLOv8 phát hiện được vùng biển số, hệ thống cắt vùng ảnh này ra và tiền xử lý. Các bước tiền xử lý giúp ảnh rõ hơn trước khi đưa vào EasyOCR. Tuy nhiên, OCR có thể nhầm một số ký tự, ví dụ O với 0, I với 1, S với 5. Vì vậy, nhóm em có thêm bước hậu xử lý dựa trên định dạng biển số Việt Nam để giảm lỗi.",
        "img": None
    },
    {
        "title": "ROLLING VOTE VÀ KẾT QUẢ MÔ HÌNH",
        "bullets": ["* Rolling Vote:", "- Lấy kết quả từ nhiều khung hình liên tiếp", "- Chọn kết quả xuất hiện ổn định nhất", "- Giảm lỗi do ảnh mờ, nghiêng, thiếu sáng", "* Kết quả đánh giá:", "- Precision = 0.979", "- Recall = 0.966", "- mAP@0.5 = 0.987"],
        "notes": "Khi nhận diện từ video, kết quả ở một khung hình đơn lẻ có thể chưa ổn định do biển số bị mờ, nghiêng hoặc thiếu sáng. Vì vậy, nhóm dùng Rolling Vote, tức là lấy kết quả từ nhiều khung hình liên tiếp và chọn kết quả ổn định nhất. Kết quả đánh giá mô hình đạt Precision 0.979, Recall 0.966 và mAP@0.5 là 0.987. Precision cao cho thấy mô hình ít nhận nhầm, Recall cao cho thấy ít bỏ sót, còn mAP cao cho thấy khả năng xác định vị trí biển số tốt.",
        "img": None
    },
    {
        "title": "PHẦN MỀM QUẢN LÝ CHO NHÂN VIÊN",
        "bullets": ["- Dashboard tổng quan", "- Theo dõi xe đang gửi", "- Xem lịch sử giao dịch", "- Quản lý người dùng", "- Tự động ghi nhận biển số khi nhận diện thành công"],
        "notes": "Phần mềm quản lý được xây dựng dành cho nhân viên bãi xe. Nhân viên có thể theo dõi xe đang gửi, xem lịch sử giao dịch, quản lý người dùng và kiểm tra thông tin liên quan đến biển số. Khi hệ thống nhận diện thành công, dữ liệu được ghi nhận tự động, giúp giảm thao tác nhập liệu thủ công.",
        "img": None
    },
    {
        "title": "APP NGƯỜI DÙNG ANDROID",
        "bullets": ["- Đăng nhập tài khoản", "- Xem số dư", "- Xem lịch sử gửi xe", "- Xem thông tin tài khoản", "- Quản lý thông tin biển số"],
        "notes": "Ở phía chủ xe, nhóm xây dựng ứng dụng Android để người dùng có thể đăng nhập, xem số dư, xem lịch sử gửi xe và kiểm tra thông tin tài khoản. Mục tiêu là giúp người dùng chủ động theo dõi thông tin cá nhân mà không cần phụ thuộc hoàn toàn vào nhân viên bãi xe.",
        "img": img_app
    },
    {
        "title": "DEMO, KẾT QUẢ VÀ HƯỚNG PHÁT TRIỂN",
        "bullets": ["* Demo workflow:", "- Nhận diện biển số -> Lưu dữ liệu -> Hiển thị trên app", "* Kết quả đạt được:", "- Pipeline nhận diện biển số bằng AI, Desktop App, Android App, API", "* Hướng phát triển:", "- Bổ sung dữ liệu thực tế, triển khai cloud, IoT/cảm biến chỗ trống, thanh toán tự động, tăng cường bảo mật"],
        "notes": "Do thời gian trình bày có giới hạn, nhóm em xin sử dụng video demo ngắn để thể hiện đầy đủ luồng hoạt động chính của hệ thống. Ở bước đầu tiên, hệ thống nhận hình ảnh từ camera hoặc video. YOLOv8 phát hiện vùng biển số, EasyOCR đọc ký tự và kết quả được hiển thị trên giao diện. Sau khi biển số được xác nhận, dữ liệu được lưu lại trên hệ thống quản lý. Cuối cùng, ở phía người dùng, app Android hiển thị các thông tin như số dư, lịch sử gửi xe và thông tin tài khoản. Tóm lại, đề tài đã xây dựng được hệ thống gồm module nhận diện biển số bằng AI, phần mềm quản lý và app người dùng. Hệ thống giúp giảm thao tác thủ công, hỗ trợ tra cứu dữ liệu và cải thiện trải nghiệm người dùng. Phần trình bày của nhóm em xin kết thúc tại đây. Nhóm em xin cảm ơn quý thầy cô đã lắng nghe.",
        "img": None
    }
]

for d in slides_data:
    add_slide_with_image(prs, 1, d["title"], d["bullets"], notes=d["notes"], img_path=d["img"])

prs.save(r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\Slide_NCKH_Official.pptx")
print("Hoàn thành!")
