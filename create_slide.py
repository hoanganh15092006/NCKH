import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Slide 1: Tiêu đề
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "HỆ THỐNG QUẢN LÝ BÃI ĐỖ XE THÔNG MINH - SMART PARKING"
    subtitle.text = "Bảo vệ Nghiên cứu Khoa học\nSinh viên thực hiện: ...\nGiảng viên hướng dẫn: ..."

    def add_slide(title_text, bullet_points):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title_text
        tf = body_shape.text_frame
        
        for i, bp in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bp
            p.font.size = Pt(22)

    # Slide 2
    add_slide("1. Đặt vấn đề (Thực trạng)", [
        "Ùn tắc: Dừng xe, lấy thẻ, trả tiền mặt làm chậm lưu thông.",
        "Rủi ro: Thẻ từ dễ mất, vé giấy dễ hỏng hoặc làm giả.",
        "Chi phí: Đòi hỏi nhân sự túc trực tại các chốt vé 24/7.",
        "=> YÊU CẦU: Tự động hóa 100%, dùng điện thoại thông minh thay thẻ cứng."
    ])

    # Slide 3
    add_slide("2. Giải pháp đột phá: 3 KHÔNG", [
        "KHÔNG chờ đợi: Nhận diện biển số tự động ngay khi xe đến.",
        "KHÔNG chạm: Barie mở tự động, người dùng không cần lấy thẻ.",
        "KHÔNG tiền mặt: Tích hợp ví điện tử và hệ thống thanh toán thông minh."
    ])

    # Slide 4
    add_slide("3. Kiến trúc hệ thống", [
        "Cụm Camera Local (PC): Mắt thần dùng AI bóc tách ký tự biển số.",
        "Backend Server (Python Flask/MySQL): Cung cấp RESTful API tốc độ cao.",
        "Mobile App (Android/Java): Điểm chạm mượt mà kết nối qua Retrofit."
    ])

    # Slide 5
    add_slide("4. Trải nghiệm & Tính năng dự phòng", [
        "Tốc độ cao: Quá trình đọc biển -> Check số dư -> Mở Barie < 3 giây.",
        "Dự phòng sự cố (App): Cho phép người dùng tự quét mã QR tại trạm để mở cổng từ xa nếu biển số xe bị bùn đất che khuất."
    ])

    # Slide 6
    add_slide("5. Tự động hóa Thanh toán (Điểm nhấn NCKH)", [
        "Luồng thanh toán: Tạo VietQR chuyển tiền trực tiếp vào tài khoản ngân hàng chủ bãi.",
        "Ứng dụng RPA: Robot tự động quét lịch sử giao dịch MB Bank 24/7.",
        "AI Giải mã CAPTCHA: Sử dụng mạng nơ-ron CNN (TensorFlow) nhận diện hình ảnh để vượt tường lửa ngân hàng.",
        "=> Kết quả: Nạp tiền vào ví App trong vài giây, miễn phí hoàn toàn."
    ])

    # Slide 7
    add_slide("6. Giao diện Ứng dụng", [
        "(Vui lòng chèn 3-4 hình ảnh chụp màn hình App vào đây)",
        "- Giao diện Trang chủ hiện đại",
        "- Màn hình Quét QR giả lập trạm",
        "- Màn hình Theo dõi trạng thái & Lịch sử"
    ])

    # Slide 8
    add_slide("7. Video Demo Hoạt động", [
        "(Vui lòng chèn Video demo hệ thống chạy thực tế vào slide này)",
        "- Nhấn nút quét QR",
        "- Hệ thống nhận lệnh và điều khiển Barie",
        "- App trừ tiền trực tiếp"
    ])

    # Slide 9
    add_slide("8. Kết quả đạt được", [
        "Hoàn thiện hệ sinh thái khép kín: Camera - Server - Mobile App.",
        "Áp dụng thành công AI trong: Nhận diện biển số và Giải mã CAPTCHA.",
        "Hoạt động ổn định với thời gian phản hồi API dưới 200ms."
    ])

    # Slide 10
    add_slide("9. Hướng phát triển & Lời Cảm ơn", [
        "Tương lai: Tích hợp cảm biến siêu âm dẫn đường xe vào chỗ trống.",
        "Em xin chân thành cảm ơn Hội đồng đã lắng nghe!",
        "Q&A - Hỏi đáp."
    ])

    prs.save("c:\\Users\\Admin\\CAMNHANDIENBIENSO\\NCKH\\Slide_BaoCao_SmartParking.pptx")
    print("Đã tạo thành công file PowerPoint!")

if __name__ == '__main__':
    create_presentation()
