import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def clear_slides(prs):
    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    for slide in slides:
        xml_slides.remove(slide)

def add_slide_with_image(prs, layout_idx, title, bullets, img_path=None):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    
    if slide.shapes.title:
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153) # Navy blue
        
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
            
    if body and bullets:
        tf = body.text_frame
        tf.clear()
        for i, text in enumerate(bullets):
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(64, 64, 64) # Dark gray
            
    if img_path:
        left = Inches(5.2)
        top = Inches(2.0)
        height = Inches(4.5)
        try:
            slide.shapes.add_picture(img_path, left, top, height=height)
        except Exception as e:
            print(f"Lỗi chèn ảnh {img_path}: {e}")

prs = Presentation(r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\slide NCKH.pptx")
clear_slides(prs)

# Title Slide
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
if title_slide.shapes.title:
    title_slide.shapes.title.text = "HỆ THỐNG QUẢN LÝ BÃI ĐỖ XE THÔNG MINH\n(SMART PARKING)"
    title_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)
try:
    title_slide.placeholders[1].text = "Bảo vệ Nghiên cứu Khoa học\nSinh viên thực hiện: Hoàng Anh\nGiảng viên hướng dẫn: [Điền tên Giảng Viên]"
except:
    pass

# Image paths
img_gate = r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\web_parking.jpg"
img_app = r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\web_app.jpg"
img_ai = r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\web_ai.jpg"

# Content
add_slide_with_image(prs, 1, "1. Thực trạng bãi đỗ xe truyền thống", [
    "- Ùn tắc nghiêm trọng tại cổng ra/vào giờ cao điểm.",
    "- Thao tác thủ công (phát thẻ, thu tiền lẻ) tốn rất nhiều thời gian.",
    "- Ô nhiễm tiếng ồn và khói bụi do xe phải dừng chờ quá lâu.",
    "- Phụ thuộc hoàn toàn vào nhân viên bảo vệ."
], img_gate)

add_slide_with_image(prs, 1, "2. Rủi ro và Bất cập của Vé giấy/Thẻ từ", [
    "- Rủi ro bảo mật: Vé giấy dễ bị làm giả, thẻ từ dễ bị đánh tráo.",
    "- Thất thoát tài chính: Quy trình thu tiền mặt khó kiểm soát minh bạch.",
    "- Chi phí vận hành lớn: Phải duy trì đội ngũ bảo vệ túc trực 24/7.",
    "- Chi phí vật tư: Thay thế thẻ từ hư hỏng, in ấn vé giấy liên tục."
])

add_slide_with_image(prs, 1, "3. Mục tiêu nghiên cứu của Đề tài", [
    "1. Số hóa quy trình: Thay thế hoàn toàn thẻ vật lý bằng công nghệ nhận diện.",
    "2. Tự động hóa thanh toán: Áp dụng Fintech để tạo vòng lặp khép kín.",
    "3. Nâng cao trải nghiệm: Giảm thiểu thời gian chờ đợi của khách hàng.",
    "4. Tối ưu chi phí: Xây dựng hệ thống tự hành, giảm phụ thuộc con người."
])

add_slide_with_image(prs, 1, "4. Giải pháp công nghệ: Triết lý 3 KHÔNG", [
    "✅ KHÔNG CHỜ ĐỢI: Nhận diện biển số tức thời bằng Computer Vision.",
    "✅ KHÔNG CHẠM: Barie tự động nhấc lên khi xe tiến vào.",
    "✅ KHÔNG TIỀN MẶT: Tích hợp ví điện tử trên Mobile App.",
    "=> Biến chiếc điện thoại thông minh thành một tấm vé đa năng."
])

add_slide_with_image(prs, 1, "5. Tổng quan Kiến trúc Hệ thống", [
    "Hệ thống phân tán (Distributed System) gồm 3 lớp chính:",
    "- Lớp Hardware (Local): Cụm Camera IP quét biển số và Barie vật lý.",
    "- Lớp Middleware (Backend): Server trung tâm xây dựng bằng Python Flask, quản trị Database MySQL.",
    "- Lớp Client (Mobile App): Ứng dụng Android viết bằng Java, giao tiếp qua Retrofit API."
])

add_slide_with_image(prs, 1, "6. Cụm nhận diện Camera Local (AI Vision)", [
    "- Hoạt động 24/7 tại các luồng xe ra vào.",
    "- Áp dụng công nghệ Optical Character Recognition (OCR).",
    "- Bóc tách, nhận diện và chuẩn hóa ký tự biển số trong vài mili-giây.",
    "- Chống nhiễu khi thời tiết xấu hoặc biển số bị mờ."
], img_gate)

add_slide_with_image(prs, 1, "7. Backend Server - Bộ não của Hệ thống", [
    "- Phát triển bằng Python Flask: Đảm bảo hiệu năng và dễ mở rộng.",
    "- Kiến trúc RESTful API: Phục vụ cả Mobile App và Cụm Camera.",
    "- Giao tiếp thời gian thực: Đồng bộ trạng thái Barie với App Mobile.",
    "- Bảo mật: Lưu trữ mật khẩu và lịch sử giao dịch an toàn trên MySQL."
])

add_slide_with_image(prs, 1, "8. Trải nghiệm Mobile App (Client)", [
    "- Giao diện thiết kế theo phong cách tối giản (Minimalist), tông màu hiện đại.",
    "- Tính năng lõi: Quản lý phương tiện cá nhân, Nạp tiền, Lịch sử ra vào.",
    "- Tính năng dự phòng sự cố: Cho phép người dùng tự quét mã QR tại trạm để mở cổng nếu Camera bị lỗi nhận diện."
], img_app)

add_slide_with_image(prs, 1, "9. Thách thức lớn nhất: Thanh toán Tự động", [
    "- Thực tế: Các bãi xe muốn tích hợp thanh toán ngân hàng thường phải trả phí rất đắt cho cổng VNPAY, MoMo.",
    "- Nếu dùng chuyển khoản thường: Admin phải ngồi kiểm tra App ngân hàng rồi cộng tiền tay, rất bất tiện.",
    "=> Câu hỏi: Làm sao để đối soát tiền vào tài khoản ngân hàng HOÀN TOÀN TỰ ĐỘNG và MIỄN PHÍ?"
])

add_slide_with_image(prs, 1, "10. Giải pháp Thanh toán: Công nghệ RPA", [
    "- RPA (Robotic Process Automation) - Tự động hóa bằng Robot.",
    "- Xây dựng một Script NodeJS sử dụng Puppeteer chạy nền trên Server.",
    "- Robot này sẽ đóng vai một 'kế toán viên', tự động đăng nhập vào trang web Internet Banking của MB Bank.",
    "- Truy xuất lịch sử giao dịch liên tục mỗi 5 giây."
])

add_slide_with_image(prs, 1, "11. Rào cản kỹ thuật: Tường lửa Ngân hàng", [
    "- Để chống lại Robot, ngân hàng MB Bank yêu cầu nhập mã CAPTCHA hình ảnh mỗi khi đăng nhập.",
    "- Các công cụ RPA thông thường sẽ bị vô hiệu hóa hoàn toàn tại bước này.",
    "- Vấn đề này đòi hỏi phải có Trí tuệ nhân tạo (AI) can thiệp để xử lý."
])

add_slide_with_image(prs, 1, "12. Đột phá Công nghệ: Mạng Nơ-ron (CNN)", [
    "- Tự tay thiết kế và huấn luyện mô hình Mạng Nơ-ron Tích chập (CNN).",
    "- Framework: TensorFlow (Google).",
    "- Dữ liệu: Hàng nghìn mẫu CAPTCHA thực tế được gán nhãn.",
    "- Khả năng: Khử nhiễu nền ảnh, cắt tách từng ký tự và dự đoán chính xác nội dung.",
    "- Giúp Robot vượt rào bảo mật ngân hàng một cách dễ dàng."
], img_ai)

add_slide_with_image(prs, 1, "13. Quy trình Thanh toán Khép kín (End-to-End)", [
    "- Bước 1: App tạo mã VietQR theo định dạng chuẩn.",
    "- Bước 2: Khách hàng quét mã chuyển khoản.",
    "- Bước 3: Robot (RPA + AI) nhận diện tiền về tài khoản ngân hàng chủ bãi.",
    "- Bước 4: Robot gửi HTTP Request lên Backend Server.",
    "- Bước 5: Server cập nhật số dư ví trong vòng 3 giây. Hệ thống hoàn toàn tự vận hành."
])

add_slide_with_image(prs, 1, "14. Đánh giá Kết quả Thực nghiệm", [
    "- Hoàn thiện toàn bộ hệ sinh thái phần mềm: Hardware - Backend - Mobile.",
    "- Tốc độ xử lý ưu việt: Thời gian trung bình một xe qua cổng giảm từ 15s xuống còn dưới 3s.",
    "- Tính ổn định: Backend xử lý mượt mà, không gặp hiện tượng thắt cổ chai dữ liệu.",
    "- Ứng dụng thành công AI vào bài toán thực tiễn phức tạp (Thanh toán ngân hàng)."
])

add_slide_with_image(prs, 1, "15. Định hướng Mở rộng trong tương lai", [
    "- Indoor Navigation: Tích hợp cảm biến siêu âm kết hợp bản đồ 2D trên App để dẫn đường đỗ xe.",
    "- Nền tảng iOS: Xây dựng ứng dụng phiên bản Swift/Flutter để hỗ trợ người dùng iPhone.",
    "- Đa dạng Ngân hàng: Cập nhật mô hình AI để đối soát tự động với các ngân hàng lớn khác."
])

add_slide_with_image(prs, 0, "XIN CHÂN THÀNH CẢM ƠN HỘI ĐỒNG!", [
    "Bài thuyết trình đến đây là kết thúc.",
    "Em rất mong nhận được những câu hỏi, phản biện và đóng góp ý kiến từ quý Thầy Cô để đề tài NCKH được hoàn thiện hơn."
])

prs.save(r"c:\Users\Admin\CAMNHANDIENBIENSO\NCKH\Slide_NCKH_ChiTiet.pptx")
print("Hoàn thành!")
